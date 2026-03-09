import json
import os
import re
import sys
import webbrowser
from concurrent.futures import ThreadPoolExecutor, as_completed
import pyautogui
import requests
import time
from PySide6.QtGui import QIcon
from mutagen import File
from pptx import Presentation
from docx import Document
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
from PySide6.QtCore import Signal, QThread, QUrl, QTimer, Qt, QSize
from PySide6.QtMultimedia import QMediaPlayer, QAudioOutput, QMediaDevices
from PySide6.QtWidgets import QFileDialog, QApplication, QWidget, QLabel, QVBoxLayout, QHBoxLayout, QLineEdit
from qfluentwidgets import (FluentWindow, FluentIcon, InfoBar, InfoBarPosition, NavigationItemPosition,
                            InfoLevel, setThemeColor, RoundMenu, Action, MessageBoxBase, SubtitleLabel,
                            LineEdit, MessageBox, SplashScreen, PushButton, CardWidget, CaptionLabel, ComboBox,
                            SpinBox, DoubleSpinBox)

from Ui_mainwindow import Ui_mainwindow
from settingInterface import Ui_settingInterface
from toolsInterface import Ui_toolsInterface
from tts_engine import TTSEngine


class PPTReviewer(QWidget, Ui_mainwindow):
    """主要实现"""

    def __init__(self, parent=None):
        super().__init__(parent=parent)
        self.setupUi(self)

        self.player = QMediaPlayer()
        self.audio_output = QAudioOutput()
        self.player.setAudioOutput(self.audio_output)
        self.player.mediaStatusChanged.connect(self.media_status_changed)
        
        # 监听音频设备物理插拔
        self._media_devices = QMediaDevices(self)
        self._media_devices.audioOutputsChanged.connect(self.handle_audio_device_change)
        
        # 轮询监听Windows系统“软”切换默认输出设备
        self.current_device_id = QMediaDevices.defaultAudioOutput().id()
        self.device_check_timer = QTimer(self)
        self.device_check_timer.timeout.connect(self.check_default_audio_device)
        self.device_check_timer.start(1000)
        
        self.media_list = []  # 存储音频文件的列表
        self.current_index = 0  # 当前播放的音频文件索引
        self.wait_current_index = 0  # 倒计时索引
        self.note_file_path = ''
        self.note_file_name = ''
        self.notes = {}  # 每页讲稿
        self.notes_list = []  # 每块讲稿
        self.notes_duration_list = [0]
        self.is_play_notes = False
        self.is_import = False
        self.mark = '●'

        self.partingIconWidget.setIcon(QIcon(':/image/image/parting.svg'))
        self.fileIconWidget.setIcon(QIcon(':/image/image/Folder.svg'))
        self.currentIconWidget.setIcon(QIcon(':/image/image/countdown.svg'))
        self.currentTimeIconWidget.setIcon(QIcon(':/image/image/Clock.svg'))
        self.pageJumpToolButton.setIcon(FluentIcon.ACCEPT_MEDIUM)
        # self.infoPushButton.setIcon(FluentIcon.IOT)
        self.playButton.setIcon(QIcon(':/image/image/play.svg'))
        self.stopButton.setIcon(QIcon(':/image/image/stop.svg'))
        self.resetButton.setIcon(QIcon(':/image/image/backward.svg'))
        self.IndeterminateProgressBar.setVisible(False)

        self.getFileButton.setIcon(QIcon(':/image/image/ppt.svg'))
        self.file_button_menu = RoundMenu(parent=self)
        self.file_button_menu.addAction(
            Action(QIcon(':/image/image/word.svg'), '导入 Word', triggered=self.init_word_play))
        self.getFileButton.setFlyout(self.file_button_menu)
        self.bgScrollArea.enableTransparentBackground()

        self.playButton.clicked.connect(self.init_play)
        self.stopButton.clicked.connect(self.stop_audio)
        self.resetButton.clicked.connect(self.reset_audio)
        self.getFileButton.clicked.connect(self.init_ppt_play)
        self.editMarkPushButton.clicked.connect(self.show_edit_mark_dialog)
        self.pageJumpToolButton.clicked.connect(self.jump_page)
        self.infoPushButton.clicked.connect(self.show_info_dialog)

        self.next_timer = QTimer(self)
        self.next_timer.timeout.connect(self.timeout_play_next_audio)  # 当计时器超时时连接到play_next_audio方法

        self.wav_temp_path = os.path.abspath('./temp')
        self.countdown_wav_temp_path = os.path.abspath('./temp/countdown')

        os.makedirs(self.wav_temp_path, exist_ok=True)
        os.makedirs(self.countdown_wav_temp_path, exist_ok=True)

        self.save_thread = SaveThread(self, tts, self)
        self.save_thread.signal_import_index.connect(self.thread_print_index)
        self.save_thread.signal_finish.connect(self.thread_save_finish)

        self.check_import()

    def check_import(self):
        """检查导入状态"""
        if self.is_import:
            self.playCardWidget.setEnabled(True)
            self.playCardWidget_2.setEnabled(True)
            self.playCardWidget_3.setEnabled(True)
            self.statusLabel.setText('已导入')
            self.IconInfoBadge.setLevel(InfoLevel.SUCCESS)
            self.IconInfoBadge.setIcon(FluentIcon.ACCEPT_MEDIUM)
        else:
            self.playCardWidget.setEnabled(False)
            self.playCardWidget_2.setEnabled(False)
            self.playCardWidget_3.setEnabled(False)
            self.statusLabel.setText('未导入')
            self.IconInfoBadge.setLevel(InfoLevel.INFOAMTION)
            self.IconInfoBadge.setIcon(FluentIcon.ACCEPT_MEDIUM)

    def handle_audio_device_change(self):
        """处理设备物理插拔的自适应切换"""
        self.check_default_audio_device()

    def check_default_audio_device(self):
        """轮询或事件触发：发现默认设备变化时自动切换"""
        default_device = QMediaDevices.defaultAudioOutput()
        if default_device.id() != self.current_device_id:
            self.current_device_id = default_device.id()
            self.audio_output.setDevice(default_device)
            # 在某些情况下强刷媒体对象可能有助于它立刻生效
            if self.player.playbackState() == QMediaPlayer.PlaybackState.PlayingState:
                pos = self.player.position()
                self.player.pause()
                self.player.play()
                self.player.setPosition(pos)
            print(f'[播放器] 音频输出设备已自适应切换至: {default_device.description()}')

    def play_audio(self):
        """播放音频文件"""
        if self.is_play_notes:  # 播放讲稿
            if self.current_index < len(self.media_list):  # 正在播放
                self.player.setSource(QUrl.fromLocalFile(self.media_list[self.current_index]))
                self.player.play()
                self.playButton.setEnabled(False)  # 禁用播放按钮
                self.currentStatusLabel.setText('播放')
                self.set_current_label_text()
                print(self.notes_list[self.current_index]["text"])
            else:
                print('播放完毕')
                self.reset_audio()

        else:  # 播放倒计时
            if self.wait_current_index < len(self.media_list):  # 正在播放
                self.player.setSource(QUrl.fromLocalFile(self.media_list[self.wait_current_index]))
                self.player.play()
                self.playButton.setEnabled(False)  # 禁用播放按钮
                temp_index = len(self.media_list) - self.wait_current_index
                self.currentStatusLabel.setText('倒计时')
                self.currentPageLabel.setText(f'{temp_index}')
                self.currentIndexLabel.setText(f'{temp_index}')
            else:
                print('播放完毕')
                self.wait_current_index = 0
                self.play_notes()

    def media_status_changed(self, status):
        """媒体状态转移"""
        # 播放结束，切换至下一个音频
        if status == QMediaPlayer.MediaStatus.EndOfMedia:
            self.next_timer.start(100)

    def stop_audio(self):
        """停止播放"""
        self.player.stop()
        self.playButton.setEnabled(True)  # 启用播放按钮
        self.currentStatusLabel.setText('停止')
        
        # 停止试听播放器
        # 尝试获取主窗口的 setting_interface
        try:
            if self.window() and hasattr(self.window(), 'setting_interface'):
                if hasattr(self.window().setting_interface, 'preview_player'):
                    self.window().setting_interface.preview_player.stop()
        except Exception:
            pass

    def reset_audio(self):
        """重置播放"""
        self.stop_audio()
        self.current_index = 0
        self.set_current_label_text()

    def set_current_label_text(self):
        self.currentPageLabel.setText(f'{self.notes_list[self.current_index]["page"]} / {len(self.notes)}')
        self.currentIndexLabel.setText(f'{self.current_index + 1} / {len(self.notes_list)}')

    def get_index_from_page(self, page):
        notes_list = self.notes_list
        for i, item in enumerate(notes_list):
            if item['page'] == page:
                return i
        return -1

    def timeout_play_next_audio(self):
        """定时结束，自动播放下一个音频"""
        self.next_timer.stop()
        if self.is_play_notes:
            if self.scrollEnableSwitch.isChecked():  # 同步翻页
                pyautogui.press("pagedown")
            self.current_index += 1
        else:
            self.wait_current_index += 1
        self.play_audio()

    def init_word_play(self):
        """初始化word"""
        self.getFileButton.setEnabled(False)
        self.get_word_path()  # 获取word文档地址
        if not self.note_file_path:  # 未选择地址，警告
            print('文件未导入！')
            self.create_warning_info_bar('导入已取消', '请重新选择文件进行导入。')
            self.getFileButton.setEnabled(True)
            return False
        self.notesPathLabel.setText(self.note_file_path)  # 设置label

        setThemeColor('#2B579A')

        try:
            self.get_word_notes_dict()  # 解析word备注
        except Exception as e:  # 解析失败，报错
            print(e)
            self.create_error_info_bar('Word 文件解析错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False
        self.init_general_play()

    def init_ppt_play(self):
        """初始化PPT"""
        self.getFileButton.setEnabled(False)

        self.get_ppt_path()
        if not self.note_file_path:  # 未选择地址，警告
            print('文件未导入！')
            self.create_warning_info_bar('导入已取消', '请重新选择文件进行导入。')
            self.getFileButton.setEnabled(True)
            return False
        self.notesPathLabel.setText(self.note_file_path)  # 设置label

        setThemeColor('#B7472A')

        try:
            self.get_ppt_notes_dict()
        except Exception as e:  # 解析失败，报错
            print(e)
            self.create_error_info_bar('PowerPoint 文件解析错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False
        self.init_general_play()
        return True

    def init_general_play(self):
        """通用初始化"""
        try:
            self.mark_split()  # 分割备注
        except Exception as e:  # 失败，报错
            print(e)
            self.create_error_info_bar('讲稿解析错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False

        self.clean_and_reset()

        try:
            self.clean_temp_folder(self.wav_temp_path)  # 清理缓存语音
            self.clean_temp_folder(self.countdown_wav_temp_path)  # 清理缓存语音
        except Exception as e:
            print(e)
            self.create_error_info_bar('缓存清理错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False

        try:
            self.save_thread.start()  # 解析成功，开始转换语音
            self.IndeterminateProgressBar.setVisible(True)
        except Exception as e:
            print(e)
            self.create_error_info_bar('语音转换错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False

    def clean_and_reset(self):
        """清理可能已导入的内容"""
        self.stop_audio()
        self.player.setSource(QUrl()) # 释放主播放器资源
        
        # 释放试听播放器资源
        try:
            if self.window() and hasattr(self.window(), 'setting_interface'):
                if hasattr(self.window().setting_interface, 'preview_player'):
                    self.window().setting_interface.preview_player.stop()
                    self.window().setting_interface.preview_player.setSource(QUrl())
        except Exception:
            pass

        self.media_list = []
        self.current_index = 0
        self.is_import = False
        self.check_import()

    def get_word_path(self):
        """获取word路径"""
        selected_files = QFileDialog.getOpenFileName(self, "选择Word文件", "", "Word Files (*.docx)")
        self.note_file_path = selected_files[0]
        self.set_filename()

    def get_ppt_path(self):
        """获取ppt路径"""
        selected_files = QFileDialog.getOpenFileName(self, "选择PowerPoint文件", "", "PowerPoint Files (*.pptx)")
        self.note_file_path = selected_files[0]
        self.set_filename()

    def set_filename(self):
        """根据文件路径生成文件名"""
        filename = os.path.basename(self.note_file_path)
        self.note_file_name = os.path.splitext(filename)[0]

    def get_word_notes_dict(self):
        """从word获取讲稿字典"""
        doc = Document(self.note_file_path)
        page_data = {}  # 用于存储页数和内容的字典
        current_page = None  # 当前页码
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            # 检查段落中是否包含“第x页”标记
            if text.startswith("第") and text.endswith("页"):
                # 提取页码
                page_text = text[1:-1]  # 去掉“第”和“页”
                try:
                    current_page = int(page_text)
                except ValueError:
                    current_page = None
            elif current_page is not None:
                if current_page in page_data:
                    page_data[current_page] += '\n' + text
                else:
                    page_data[current_page] = text
        self.notes = page_data

    def get_ppt_notes_dict(self):
        """从PPT备注获取讲稿字典"""
        presentation = Presentation(self.note_file_path)
        i = 1
        notes_dict = {}
        # 遍历 PPT 中的每一页
        for slide in presentation.slides:
            # 遍历每一页中的每一个备注
            text = ''
            for note in slide.notes_slide.notes_text_frame.paragraphs:
                text += note.text + '\n'
            notes_dict[i] = text
            i += 1
        self.notes = notes_dict

    def mark_split(self):
        """按标记符分割讲稿"""
        notes_list = []
        for page in range(1, len(self.notes) + 1):
            note_text = re.sub(r'[\x00-\x1f\x7f]+', '', self.notes[page])
            if note_text:
                if note_text[-1] == self.mark:
                    note_text = note_text[:-1]
            note_list = note_text.split(self.mark)  # 按分隔符分割
            for one_note in note_list:
                notes_list.append({'page': page, 'text': one_note.strip()})
        self.notes_list = notes_list
        print('讲稿分割完毕')

    @staticmethod
    def clean_temp_folder(path):
        """清理缓存wav"""
        for filename in os.listdir(path):
            if filename.endswith(".wav"):
                file_path = os.path.join(path, filename)
                try:
                    os.remove(file_path)
                    print(f"已清理: {filename}")
                except Exception as e:
                    print(f"清理文件失败: {filename}, 原因: {e}")
        print('清理完成')

    def thread_print_index(self, import_index):
        """信号，生成数据更新"""
        text = f'已生成：{import_index}/{len(self.notes_list)}'
        print(text)
        self.statusLabel.setText(text)

    def thread_save_finish(self):
        """信号，转换完成"""
        print('转换完成！')
        self.IndeterminateProgressBar.setVisible(False)
        self.create_success_info_bar('转换完成', '语音播放功能已准备就绪。')
        self.getFileButton.setEnabled(True)
        self.is_import = True
        self.check_import()
        self.set_current_label_text()
        self.pageJumpSpinBox.setMaximum(len(self.notes))

    def init_play(self):
        """准备播放"""
        if self.currentSwitch.isChecked():
            # 载入倒计时，播放
            self.play_wait()
        else:
            # 载入音频，播放
            self.play_notes()

    def play_wait(self):
        """播放倒计时"""
        self.is_play_notes = False
        self.wait_current_index = 0
        self.load_wait_audio_files()
        print('已导入倒计时')
        self.play_audio()

    def play_notes(self):
        """播放讲稿"""
        self.is_play_notes = True
        if not self.media_list or 'countdown' in self.media_list[0]:  # 还未载入，必须载入；载入的为倒计时，需重新载入
            self.load_audio_files()
        self.play_audio()

    def load_audio_files(self):
        """查找所有.wav，添加到media_list中"""
        audio_files = [f for f in os.listdir(self.wav_temp_path) if f.endswith(".wav")]
        audio_files = sorted(audio_files, key=lambda x: [int(part) if part.isdigit() else part
                                                         for part in os.path.splitext(x)[0].split("_")])
        self.media_list = [f'{self.wav_temp_path}/{f}' for f in audio_files]
        print('音频列表载入完成')

    def load_wait_audio_files(self):
        """获取指定数量的倒计时wav，添加到media_list中"""
        audio_files = [f for f in os.listdir(self.countdown_wav_temp_path)
                       if f.endswith(".wav") and int(f.split('.')[0]) <= self.currentSpinBox.value()]
        audio_files = sorted(audio_files, key=lambda x: int(os.path.splitext(x)[0]), reverse=True)
        self.media_list = [f'{self.countdown_wav_temp_path}/{f}' for f in audio_files]
        print('倒计时列表载入完成')

    def jump_page(self):
        self.stop_audio()
        index = self.get_index_from_page(self.pageJumpSpinBox.value())
        if index > -1:
            self.current_index = index
        self.set_current_label_text()

    def create_success_info_bar(self, title, text):
        """成功消息框"""
        InfoBar.success(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=2000,
            parent=self
        )

    def create_warning_info_bar(self, title, text):
        """警告消息框"""
        InfoBar.warning(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=3000,
            parent=self
        )

    def create_error_info_bar(self, title, text):
        """错误消息框"""
        InfoBar.error(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=-1,
            parent=self
        )

    def show_edit_mark_dialog(self):
        """编辑分隔符弹窗"""
        box = EditMarkMessageBox(self)
        if box.exec():
            text = box.urlLineEdit.text()
            print(f'分隔符：{text}')
            self.mark = text

    @staticmethod
    def s_to_str(s):
        """将秒转换为时间标签"""
        if s < 60:
            return f'{round(s, 2)} 秒'
        m, s = divmod(s, 60)
        if m < 60:
            return f'{round(m)} 分钟 {round(s)} 秒'
        else:
            h, m = divmod(m, 60)
            return f'{round(h)} 小时 {round(m)} 分钟 {round(s)} 秒'

    def count_words(self):
        """字数统计"""
        text = ''
        for page in self.notes:
            text += self.notes[page]
        text = re.sub(r'\s+', '', text)
        return len(text)

    def show_info_dialog(self):
        """显示统计弹窗"""
        title = '统计信息'
        words_count = self.count_words()
        max_duration = max(self.notes_duration_list)
        max_duration_index = self.notes_duration_list.index(max_duration)
        min_duration = min(self.notes_duration_list)
        min_duration_index = self.notes_duration_list.index(min_duration)
        content_list = [
            ['页码总计', f'{len(self.notes)} 页'],
            ['音频总计', f'{len(self.notes_list)} 条'],
            ['演讲稿字数总计', f'{words_count} 字'],
            ['音频总时长', f'{self.s_to_str(sum(self.notes_duration_list))}\n'],
            ['最长音频时长', f'{self.s_to_str(max_duration)}'],
            ['最长音频索引', f'{max_duration_index}'],
            ['最长音频所属页码', f'第 {self.notes_list[max_duration_index]["page"]} 页'],
            ['最短音频时长', f'{self.s_to_str(min_duration)}'],
            ['最短音频索引', f'{min_duration_index}'],
            ['最短音频所属页码', f'第 {self.notes_list[min_duration_index]["page"]} 页']
        ]

        content = ''
        for item in content_list:
            content += f'{item[0]}：{item[1]}\n'
        dialog = MessageBox(title, content, self)
        dialog.yesButton.setText('确定')
        dialog.cancelButton.setVisible(False)
        dialog.exec()


class SaveThread(QThread):
    """生成语音线程"""
    signal_import_index = Signal(int)
    signal_finish = Signal()

    def __init__(self, ppt_r, tts_engine: TTSEngine, parent=None):
        super().__init__(parent)
        self.ppt_r = ppt_r
        self.tts = tts_engine

    def run(self):
        self.save_countdown_wav()
        self.save_wav()
        self.signal_finish.emit()

    def _save_one_note_wav(self, index, note_dict):
        """保存单条讲稿并返回索引及音频时长。"""
        path = f'{self.ppt_r.wav_temp_path}/{note_dict["page"]}_{index + 1}.wav'
        self.tts.save_file(note_dict['text'], path)
        audio = File(path)
        duration = audio.info.length if audio and audio.info else 0
        return index, duration

    def save_wav(self):
        """调用 TTS 保存文字为 wav（支持在线引擎并行提速、自动重试）。"""
        notes_list = self.ppt_r.notes_list
        total = len(notes_list)
        info_list = [0.0] * total

        # 在线引擎默认并行，本地引擎默认串行（稳定优先）
        if self.tts.can_parallel_generate() and total > 1:
            max_workers = max(1, min(self.tts.get_parallel_workers(), total))
            completed = 0
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                future_map = {
                    executor.submit(self._save_one_note_wav, index, note_dict): index
                    for index, note_dict in enumerate(notes_list)
                }

                for future in as_completed(future_map):
                    index = future_map[future]
                    try:
                        result_index, duration = future.result()
                        info_list[result_index] = duration
                    except Exception as e:
                        raise RuntimeError(f'并行生成失败，索引 {index + 1}：{e}') from e

                    completed += 1
                    self.signal_import_index.emit(completed)
        else:
            for index, note_dict in enumerate(notes_list):
                result_index, duration = self._save_one_note_wav(index, note_dict)
                info_list[result_index] = duration
                self.signal_import_index.emit(index + 1)

        self.ppt_r.notes_duration_list = info_list

    def save_countdown_wav(self):
        """生成倒计时语音"""
        for index, time_num in enumerate(range(self.ppt_r.currentSpinBox.maximum(), 0, -1)):
            self.tts.save_file_for_stable_local(f'{time_num}', f'{self.ppt_r.countdown_wav_temp_path}/{time_num}.wav')
        print('倒计时生成完成')


class UpdateThread(QThread):
    """检查更新线程"""
    signal_finish = Signal(list)

    def run(self):
        self.get_update()

    def get_update(self):
        """HTTP从gitee获取更新"""
        try:
            response = requests.get("https://gitee.com/api/v5/repos/pth2000/PowerPointReviewer/releases/latest")
        except Exception as e:
            print(e)
            self.signal_finish.emit([2, '获取更新失败', f'详情：{e}'])
            return False

        try:
            res = response.json()
            latest_version = res['tag_name']
            latest_version = latest_version[1:]
            latest_version_name = res['name']
            latest_version_time = res['created_at']
            latest_version_download_url = res['assets'][0]['browser_download_url']

            if self.compare_versions(VERSION, latest_version):
                str_update = f'当前版本为最新版\n服务器版本：{latest_version}\n更新时间：{latest_version_time}'
                self.signal_finish.emit([0, '获取更新成功', str_update])
            else:
                str_update = (f"发现新版本！{VERSION} --> {latest_version}\n"
                              f"更新内容：{latest_version_name}\n更新时间：{latest_version_time}")
                self.signal_finish.emit([1, '获取更新成功', str_update, latest_version_download_url])
        except Exception as e:
            print(e)
            self.signal_finish.emit([2, 'HTTP 解析失败', f'详情：{e}'])
            return False

    @staticmethod
    def compare_versions(version1, version2):
        """版本号比较，1>=2 -> True"""
        parts1 = version1.split('.')
        parts2 = version2.split('.')

        # 找到较短版本号的长度
        min_length = min(len(parts1), len(parts2))

        # 逐一比较各个部分
        for i in range(min_length):
            if int(parts1[i]) < int(parts2[i]):
                return False  # version1 < version2
            elif int(parts1[i]) > int(parts2[i]):
                return True  # version1 > version2

        # 如果迄今为止所有部分都相等，检查较长版本号的余下部分
        if len(parts1) < len(parts2):
            return False  # version1 < version2
        elif len(parts1) > len(parts2):
            return True  # version1 > version2

        return True  # 两个版本号相等


class EditMarkMessageBox(MessageBoxBase):
    """编辑分隔符界面"""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.titleLabel = SubtitleLabel('编辑分隔符', self)
        self.urlLineEdit = LineEdit(self)

        self.urlLineEdit.setPlaceholderText('请输入您的讲稿分隔符，如“●”')
        self.urlLineEdit.setClearButtonEnabled(True)
        self.urlLineEdit.setText(w.ppt_r.mark)

        self.viewLayout.addWidget(self.titleLabel)
        self.viewLayout.addWidget(self.urlLineEdit)

        self.yesButton.setText('保存')
        self.cancelButton.setText('取消')

        self.widget.setMinimumWidth(350)
        if not self.urlLineEdit.text():
            self.yesButton.setDisabled(True)
        self.urlLineEdit.textChanged.connect(self._validate_url)

    def _validate_url(self, text):
        self.yesButton.setEnabled(QUrl(text).isValid())


class ToolsInterface(QWidget, Ui_toolsInterface):
    """工具页"""

    def __init__(self, parent=None):
        super().__init__(parent=parent)
        self.setupUi(self)
        self.bgScrollArea.enableTransparentBackground()

        self.transformPPTPushButton.clicked.connect(self.write_to_ppt)
        self.transformWordPushButton.clicked.connect(self.write_to_word)
        self.transformJsonPushButton.clicked.connect(self.write_to_json)

    def write_to_ppt(self):
        """导入PPT备注"""
        notes_dict = w.ppt_r.notes

        if not notes_dict:
            self.create_warning_info_bar('演讲稿未导入', '请先导入演讲稿。')
            return False

        ppt_path = self.get_ppt_path()
        if not ppt_path:
            self.create_warning_info_bar('文件选择已取消', '请重新选择文件进行导入。')
            return False

        dir_path = self.get_dir_path()
        if not dir_path:
            self.create_warning_info_bar('目录选择已取消', '请重新选择保存目录。')
            return False

        try:
            ppt = Presentation(ppt_path)
            for slide in ppt.slides:
                slide_number = ppt.slides.index(slide) + 1  # PowerPoint中的幻灯片索引从0开始，所以需要加1
                if slide_number in notes_dict:
                    notes_slide = slide.notes_slide
                    notes_text = notes_dict[slide_number]
                    notes_slide.notes_text_frame.text = notes_text
            ppt_name = ppt_path.split('/')[-1][:-5]
        except Exception as e:
            print(e)
            self.create_error_info_bar('PowerPoint 读取错误', f'详情：{e}')
            return False

        path = f'{dir_path}/{ppt_name}_NEW.pptx'

        try:
            ppt.save(path)
        except Exception as e:
            print(e)
            self.create_error_info_bar('PowerPoint 保存错误', f'详情：{e}')
            return False
        self.create_success_info_bar('生成成功', f'讲稿备注已导入完毕，文件路径：{path}')

    def write_to_word(self):
        """写入word"""
        notes = w.ppt_r.notes
        file_name = w.ppt_r.note_file_name

        if not notes:
            self.create_warning_info_bar('演讲稿未导入', '请先导入演讲稿。')
            return False

        dir_path = self.get_dir_path()
        file_name = f'{file_name}_Notes.docx'
        file_path = f'{dir_path}/{file_name}'
        if not dir_path:
            self.create_warning_info_bar('目录选择已取消', '请重新选择保存目录。')
            return False

        try:
            document = Document()
            for page_number, page_note in notes.items():
                # 写入页面编号
                p_page_number = document.add_paragraph()
                p_page_number.style = 'Heading 1'
                p_page_number.paragraph_format.space_before = Pt(0)  # 段前
                p_page_number.paragraph_format.space_after = Pt(0)  # 段后
                p_page_number.paragraph_format.line_spacing = 1.5
                run_page_number = p_page_number.add_run(f"第{page_number}页")
                run_page_number.font.name = 'Microsoft YaHei'
                run_page_number.font.size = Pt(14)  # 设置字号
                run_page_number.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
                run_page_number.bold = True
                run_page_number.font.color.rgb = RGBColor(0, 0, 0)
                # 写入页面备注
                p_page_note = document.add_paragraph()
                p_page_note.paragraph_format.space_before = Pt(0)  # 段前
                p_page_note.paragraph_format.space_after = Pt(0)  # 段后
                p_page_note.paragraph_format.line_spacing = 1.5
                page_note = re.sub(u"[\\x00-\\x08\\x0b\\x0e-\\x1f\\x7f]", "", page_note)
                run_page_note = p_page_note.add_run(page_note.strip("\n"))
                run_page_note.font.name = 'Microsoft YaHei'
                run_page_note.font.size = Pt(11)  # 设置字号
                run_page_note.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')
            document.save(file_path)
        except Exception as e:
            print(e)
            self.create_error_info_bar('Word 保存错误', f'详情：{e}')
            return False
        self.create_success_info_bar('转换成功', f'讲稿备注已转换完毕，文件路径：{file_path}')

    def write_to_json(self):
        """写入json"""
        notes_dict = w.ppt_r.notes
        file_name = w.ppt_r.note_file_name
        if not notes_dict:
            self.create_warning_info_bar('演讲稿未导入', '请先导入演讲稿。')
            return False

        dir_path = self.get_dir_path()
        file_name = f'{file_name}_Notes.json'
        file_path = f'{dir_path}/{file_name}'
        if not dir_path:
            self.create_warning_info_bar('目录选择已取消', '请重新选择保存目录。')
            return False

        try:
            with open(file_path, 'w') as json_file:
                json.dump(notes_dict, json_file)
        except Exception as e:
            print(e)
            self.create_error_info_bar('JSON 保存错误', f'详情：{e}')
            return False
        self.create_success_info_bar('转换成功', f'讲稿备注已转换完毕，文件路径：{file_path}')

    def get_dir_path(self):
        """获取目录"""
        return QFileDialog.getExistingDirectory(self, '选择保存文件夹', '', QFileDialog.Option.ShowDirsOnly)

    def get_ppt_path(self):
        """获取ppt路径"""
        selected_files = QFileDialog.getOpenFileName(self, "选择PowerPoint文件", "", "PowerPoint Files (*.pptx)")
        return selected_files[0]

    def create_success_info_bar(self, title, text):
        """成功消息框"""
        InfoBar.success(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=2000,
            parent=self
        )

    def create_warning_info_bar(self, title, text):
        """警告消息框"""
        InfoBar.warning(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=3000,
            parent=self
        )

    def create_error_info_bar(self, title, text):
        """错误消息框"""
        InfoBar.error(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=-1,
            parent=self
        )


class PreviewThread(QThread):
    """试听语音生成线程，防止阻塞 UI"""
    signal_finish = Signal(str)
    signal_error = Signal(str)

    def __init__(self, tts_engine, text, path, parent=None):
        super().__init__(parent)
        self.tts = tts_engine
        self.text = text
        self.path = path

    def run(self):
        try:
            self.tts.save_file(self.text, self.path)
            self.signal_finish.emit(self.path)
        except Exception as e:
            self.signal_error.emit(str(e))


class SettingInterface(QWidget, Ui_settingInterface):
    """设置页"""

    def __init__(self, parent=None):
        super().__init__(parent=parent)
        self.setupUi(self)

        # 持久化配置文件路径
        self.settings_file_path = os.path.abspath('./config/settings.json')
        os.makedirs(os.path.dirname(self.settings_file_path), exist_ok=True)

        self.versionIconWidget.setIcon(QIcon(':/image/image/update.svg'))
        self.copyrightIconWidget.setIcon(QIcon(':/image/image/info.svg'))
        self.githubButton.setIcon(FluentIcon.GITHUB)
        self.giteeButton.setIcon(QIcon(':/image/image/gitee.svg'))
        self.versionLabel.setText(VERSION)
        self.bgScrollArea.enableTransparentBackground()

        # 运行期动态控件缓存：key -> widget
        self.dynamic_option_widgets = {}
        self.dynamic_option_cards = []
        # 动态创建的音色卡片（独立追踪，随模型切换重建）
        self._voice_card = None
        self._voice_combo = None

        self.savePushButton.clicked.connect(self.save_setting)
        self.versionPrimaryPushButton.clicked.connect(self.get_update)
        self.githubButton.clicked.connect(self.open_github_url)
        self.giteeButton.clicked.connect(self.open_gitee_url)

        self.engineSelectComboBox.clear()
        self.engineSelectComboBox.addItems(tts.get_engine_names())
        self.engineSelectComboBox.currentIndexChanged.connect(self.change_tts_engine)

        # 让引擎选择卡的 info 区随窗口横向伸展
        self.horizontalLayout_14.setStretch(0, 1)
        self.engineSelectCaptionLabel.setWordWrap(True)

        # 启动时自动加载配置；若无配置则使用默认引擎
        self.load_persistent_settings()

        # 同步引擎下拉框与当前引擎索引
        self.engineSelectComboBox.blockSignals(True)
        self.engineSelectComboBox.setCurrentIndex(tts.get_mode_index())
        self.engineSelectComboBox.blockSignals(False)

        # 初始化可配置项 UI
        self.apply_engine_schema_to_ui()
        self.setup_voices_list()

        self.update_thread = UpdateThread()
        self.update_thread.signal_finish.connect(self.thread_get_update_finish)

        # 初始化试听功能
        self.preview_player = QMediaPlayer()
        self.preview_audio_output = QAudioOutput()
        self.preview_player.setAudioOutput(self.preview_audio_output)
        
        # 监听音频设备物理插拔
        self._media_devices = QMediaDevices(self)
        self._media_devices.audioOutputsChanged.connect(self.handle_audio_device_change)
        
        # 轮询监听系统级“软”切换设备
        self.current_preview_device_id = QMediaDevices.defaultAudioOutput().id()
        self.preview_device_check_timer = QTimer(self)
        self.preview_device_check_timer.timeout.connect(self.check_default_audio_device)
        self.preview_device_check_timer.start(1000)
        
        # 连接previewButton信号
        self.previewButton.clicked.connect(self.preview_audio)

    def handle_audio_device_change(self):
        """处理设备物理插拔的自适应切换"""
        self.check_default_audio_device()

    def check_default_audio_device(self):
        """发现默认设备变化时切换"""
        default_device = QMediaDevices.defaultAudioOutput()
        if default_device.id() != self.current_preview_device_id:
            self.current_preview_device_id = default_device.id()
            self.preview_audio_output.setDevice(default_device)
            # 如果正在播放试听，平滑过渡
            if self.preview_player.playbackState() == QMediaPlayer.PlaybackState.PlayingState:
                pos = self.preview_player.position()
                self.preview_player.pause()
                self.preview_player.play()
                self.preview_player.setPosition(pos)
            print(f'[试听器] 音频输出设备已自适应切换至: {default_device.description()}')

    def preview_audio(self):
        """试听当前设置"""
        try:
            # 先应用当前 UI 配置，确保试听使用的就是当前界面参数
            self.apply_ui_settings_to_engine()

            # 停止之前的播放并释放文件占用
            if self.preview_player.playbackState() == QMediaPlayer.PlaybackState.PlayingState:
                self.preview_player.stop()
            self.preview_player.setSource(QUrl())

            print(f"Main Preview Debug: Mode={tts.get_mode()}, VoiceIndex={tts.get_selected_voice_index()}")

            text = "这是一个试听音频，用于测试当前的语音设置。"
            
            # 生成带时间戳的文件名，防止多媒体占用冲突
            preview_path = os.path.abspath(f"./temp/preview_{int(time.time())}.wav")
            
            # 确保 temp 目录存在
            os.makedirs("./temp", exist_ok=True)
            
            # 禁用试听按钮以防重复点击阻塞
            self.previewButton.setEnabled(False)
            self.previewButton.setText("正在生成...")

            # 启动多线程生成试听音频
            self.preview_thread = PreviewThread(tts, text, preview_path, self)
            self.preview_thread.signal_finish.connect(self.on_preview_generated)
            self.preview_thread.signal_error.connect(self.on_preview_error)
            self.preview_thread.start()
            
        except Exception as e:
            print(e)
            self.create_error_info_bar('试听启动失败', f'详情：{e}')
            self.previewButton.setEnabled(True)
            self.previewButton.setText("试听")

    def on_preview_generated(self, path):
        """试听音频生成完毕的回调"""
        self.previewButton.setEnabled(True)
        self.previewButton.setText("试听")
        self.preview_player.setSource(QUrl.fromLocalFile(path))
        self.preview_player.play()

    def on_preview_error(self, err_msg):
        """试听音频生成失败的回调"""
        self.previewButton.setEnabled(True)
        self.previewButton.setText("试听")
        self.create_error_info_bar('试听生成失败', f'详情：{err_msg}')

    def clear_dynamic_option_cards(self):
        """清理运行期动态创建的引擎配置卡片（含音色卡）。"""
        for card in self.dynamic_option_cards:
            self.verticalLayout_2.removeWidget(card)
            card.deleteLater()
        self.dynamic_option_cards = []
        self.dynamic_option_widgets = {}
        # 一起清除音色卡
        if self._voice_card is not None:
            self.verticalLayout_2.removeWidget(self._voice_card)
            self._voice_card.deleteLater()
            self._voice_card = None
            self._voice_combo = None

    def create_dynamic_option_card(self, option_schema):
        """根据 schema 创建一张配置卡片（用于非内置项，如 pitch）。"""
        card = CardWidget(self.importWidget)
        h_layout = QHBoxLayout(card)
        h_layout.setContentsMargins(20, 20, 20, 20)

        info_layout = QVBoxLayout()
        info_layout.setSpacing(0)

        title_label = QLabel(card)
        title_label.setText(option_schema.get('label', option_schema.get('key', '配置项')))
        title_label.setStyleSheet("font: 700 10pt 'Segoe UI', 'Microsoft YaHei', 'PingFang SC';")
        info_layout.addWidget(title_label)

        desc_label = CaptionLabel(card)
        desc_label.setText(option_schema.get('description', ''))
        desc_label.setWordWrap(True)
        info_layout.addWidget(desc_label)

        h_layout.addLayout(info_layout, stretch=1)
        h_layout.addSpacing(20)

        option_type = option_schema.get('type')
        default_value = option_schema.get('default')

        if option_type == 'int':
            control = SpinBox(card)
            control.setMinimumSize(QSize(180, 33))
            control.setMaximumSize(QSize(180, 33))
            control.setRange(int(option_schema.get('min', -9999)), int(option_schema.get('max', 9999)))
            control.setSingleStep(int(option_schema.get('step', 1)))
            control.setValue(int(default_value if default_value is not None else 0))
        elif option_type == 'float':
            control = DoubleSpinBox(card)
            control.setMinimumSize(QSize(180, 33))
            control.setMaximumSize(QSize(180, 33))
            control.setRange(float(option_schema.get('min', -9999.0)), float(option_schema.get('max', 9999.0)))
            control.setSingleStep(float(option_schema.get('step', 0.1)))
            control.setValue(float(default_value if default_value is not None else 0.0))
        elif option_type == 'choice':
            control = ComboBox(card)
            control.setMinimumSize(QSize(240, 0))
            control.setMaximumSize(QSize(240, 16777215))
            choices = option_schema.get('choices', [])
            control.addItems([str(item) for item in choices])
            if choices and default_value in choices:
                control.setCurrentIndex(choices.index(default_value))
            elif choices:
                control.setCurrentIndex(0)
        elif option_type == 'password':
            control = LineEdit(card)
            control.setMinimumSize(QSize(240, 33))
            control.setEchoMode(QLineEdit.EchoMode.Password)
            control.setText(str(default_value if default_value is not None else ''))
        else:
            control = LineEdit(card)
            control.setMinimumSize(QSize(240, 33))
            control.setText(str(default_value if default_value is not None else ''))

        h_layout.addWidget(control)

        # 插入到“保存/试听”卡片之前，保证布局顺序自然
        insert_index = max(self.verticalLayout_2.indexOf(self.CardWidget_4), 0)
        self.verticalLayout_2.insertWidget(insert_index, card)

        option_key = option_schema.get('key')
        if option_key:
            self.dynamic_option_widgets[option_key] = control
        self.dynamic_option_cards.append(card)

    def apply_engine_schema_to_ui(self):
        """按当前引擎 schema 动态渲染设置项（全部动态化，按 schema 顺序生成）。"""
        self.clear_dynamic_option_cards()
        current_values = tts.get_current_option_values()
        schema_list = tts.get_current_options_schema()

        # 显示当前引擎的描述
        engine_def = tts.get_current_engine_definition()
        self.engineSelectCaptionLabel.setText(engine_def.get('description', ''))

        for item in schema_list:
            key = item.get('key')
            self.create_dynamic_option_card(item)
            control = self.dynamic_option_widgets.get(key)
            if control is None:
                continue
            # 用当前已保存值覆盖默认值
            value = current_values.get(key, item.get('default'))
            if isinstance(control, SpinBox):
                control.setValue(int(value))
            elif isinstance(control, DoubleSpinBox):
                control.setValue(float(value))
            elif isinstance(control, ComboBox):
                idx = control.findText(str(value))
                if idx >= 0:
                    control.setCurrentIndex(idx)
            elif isinstance(control, LineEdit):
                control.setText(str(value))

        # 若当前引擎有 'model' 配置项，监听变化以实时刷新音色列表
        model_widget = self.dynamic_option_widgets.get('model')
        if isinstance(model_widget, ComboBox):
            model_widget.currentTextChanged.connect(self._on_model_changed)

    def apply_ui_settings_to_engine(self):
        """读取当前 UI 控件值并写回当前引擎配置。"""
        option_values = {}
        for key, control in self.dynamic_option_widgets.items():
            if isinstance(control, SpinBox):
                option_values[key] = control.value()
            elif isinstance(control, DoubleSpinBox):
                option_values[key] = control.value()
            elif isinstance(control, ComboBox):
                option_values[key] = control.currentText()
            elif isinstance(control, LineEdit):
                option_values[key] = control.text()

        tts.apply_current_options(option_values)
        if self._voice_combo is not None:
            tts.set_voice(self._voice_combo.currentIndex())

    def setup_voices_list(self):
        """按当前引擎（及当前模型）重建发音人选择卡片，恢复上次选择。"""
        # 移除旧音色卡
        if self._voice_card is not None:
            self.verticalLayout_2.removeWidget(self._voice_card)
            self._voice_card.deleteLater()
            self._voice_card = None
            self._voice_combo = None

        voices_list = tts.get_voices_list()
        if not voices_list:
            return

        # 动态创建音色卡
        card = CardWidget(self.importWidget)
        h_layout = QHBoxLayout(card)
        h_layout.setContentsMargins(20, 20, 20, 20)

        info_layout = QVBoxLayout()
        info_layout.setSpacing(0)
        title_lbl = QLabel(card)
        title_lbl.setText('发音人选择')
        title_lbl.setStyleSheet("font: 700 10pt 'Segoe UI', 'Microsoft YaHei', 'PingFang SC';")
        info_layout.addWidget(title_lbl)
        cap_lbl = CaptionLabel(card)
        cap_lbl.setText('选择当前引擎的发音人')
        cap_lbl.setWordWrap(True)
        info_layout.addWidget(cap_lbl)

        h_layout.addLayout(info_layout, stretch=1)
        h_layout.addSpacing(20)

        combo = ComboBox(card)
        combo.setMinimumSize(QSize(240, 0))
        combo.setMaximumSize(QSize(240, 16777215))
        combo.addItems(voices_list)
        h_layout.addWidget(combo)

        # 插入到保存卡之前
        insert_idx = self.verticalLayout_2.indexOf(self.CardWidget_4)
        self.verticalLayout_2.insertWidget(insert_idx, card)

        saved_index = tts.get_selected_voice_index()
        if 0 <= saved_index < combo.count():
            combo.setCurrentIndex(saved_index)

        self._voice_card = card
        self._voice_combo = combo

    def change_tts_engine(self):
        """切换引擎：切换后按该引擎 schema 重绘可调节项。"""
        tts.set_mode(self.engineSelectComboBox.currentIndex())
        self.apply_engine_schema_to_ui()
        self.setup_voices_list()

    def _on_model_changed(self, model_name: str):
        """百炼模型切换时同步引擎配置并刷新对应音色列表。"""
        tts.set_current_option('model', model_name)
        self.setup_voices_list()

    def load_persistent_settings(self):
        """从磁盘读取配置并恢复到 TTS 管理器。"""
        if not os.path.exists(self.settings_file_path):
            return False
        try:
            with open(self.settings_file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            ok = tts.import_persistent_state(data)
            print(f'[设置] 已加载持久化配置：{self.settings_file_path}')
            return ok
        except Exception as e:
            print(f'[设置] 读取持久化配置失败：{e}')
            return False

    def save_persistent_settings(self):
        """将当前配置写入磁盘。"""
        data = tts.export_persistent_state()
        with open(self.settings_file_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        print(f'[设置] 已保存持久化配置：{self.settings_file_path}')

    def save_setting(self):
        try:
            self.apply_ui_settings_to_engine()
            self.save_persistent_settings()
        except Exception as e:
            print(e)
            self.create_error_info_bar('保存设置错误', f'详情{e}')
        else:
            self.create_success_info_bar('保存设置成功', '重新导入文件后生效。')

    def get_update(self):
        """获取版本更新"""
        self.versionPrimaryPushButton.setEnabled(False)
        self.update_thread.start()

    def thread_get_update_finish(self, data_list):
        self.versionPrimaryPushButton.setEnabled(True)
        if data_list[0] == 0:
            self.create_success_info_bar(data_list[1], data_list[2])
        elif data_list[0] == 1:
            self.show_update_dialog(data_list[1], data_list[2], data_list[3])
        else:
            self.create_error_info_bar(data_list[1], data_list[2])

    @staticmethod
    def open_github_url():
        webbrowser.open('https://github.com/pth2000')

    @staticmethod
    def open_gitee_url():
        webbrowser.open('https://gitee.com/pth2000')

    @staticmethod
    def open_update_url(url):
        webbrowser.open(url)

    def create_success_info_bar(self, title, text):
        """成功消息框"""
        InfoBar.success(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=5000,
            parent=self
        )

    def create_warning_info_bar(self, title, text):
        """警告消息框"""
        InfoBar.warning(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=5000,
            parent=self
        )

    def create_error_info_bar(self, title, text):
        """错误消息框"""
        InfoBar.error(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=-1,
            parent=self
        )

    def show_update_dialog(self, title, content, url):
        dialog = MessageBox(title, content, self)
        dialog.yesButton.setText('获取更新')
        dialog.cancelButton.setText('取消')
        if dialog.exec():
            self.open_update_url(url)


class Window(FluentWindow):
    """主窗体"""

    def __init__(self):
        super().__init__()
        setThemeColor('#B7472A')
        self.resize(850, 750)
        # 设置标题
        self.setWindowTitle('PowerPointReviewer')
        self.setWindowIcon(QIcon(':/image/image/ppt_ico.svg'))
        # 设置启动页
        self.splashScreen = SplashScreen(self.windowIcon(), self)
        self.splashScreen.setIconSize(QSize(106, 106))
        self.show()
        # 加载页面
        self.ppt_r = PPTReviewer(self)
        self.setting_interface = SettingInterface(self)
        self.tools_interface = ToolsInterface(self)
        self.addSubInterface(self.ppt_r, FluentIcon.HOME, '主页')
        self.addSubInterface(self.tools_interface, FluentIcon.APPLICATION, '实用工具')
        self.addSubInterface(self.setting_interface, FluentIcon.SETTING, '设置', NavigationItemPosition.BOTTOM)
        # 隐藏启动页
        self.splashScreen.finish()

    def click_test(self):
        self.ppt_r.create_success_info_bar('稍安勿躁', '功能还在紧锣密鼓地开发中……')


if __name__ == '__main__':
    VERSION = '1.2.0'
    tts = TTSEngine()
    app = QApplication(sys.argv)  # 声明应用程序
    w = Window()  # 声明窗口
    w.show()
    sys.exit(app.exec())
