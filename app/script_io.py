"""将 PowerPoint、Word 和 JSON 讲稿解析为统一的逐页数据结构。

PowerPoint 以幻灯片作为页边界；Word 优先解析表格，无表格时兼容手写的段落式文档；
JSON 以显式字段保留分页和页内分段。

Markdown 只承担可读导出，不参与回导。
"""

import json
import re
from dataclasses import dataclass, field
from pathlib import Path


SCRIPT_SUFFIXES = ('.pptx', '.docx', '.json')
SCRIPT_FILTER = '讲稿文件 (*.pptx *.docx *.json)'

JSON_FORMAT = 'powerpointreviewer-script'
JSON_VERSION = 2

_CN_DIGITS = {'零': 0, '〇': 0, '一': 1, '二': 2, '两': 2, '三': 3, '四': 4,
              '五': 5, '六': 6, '七': 7, '八': 8, '九': 9}


@dataclass
class ScriptData:
    """保存解析后的逐页正文、来源信息、分隔符和诊断消息。"""

    notes: dict = field(default_factory=dict)
    source_name: str = ''
    mark: str = ''          # 空串表示源文件未指定分隔符，应沿用当前应用设置。
    report: list = field(default_factory=list)


# 页码解析

def parse_chinese_number(text: str):
    """解析零至九十九的常见中文数字，无法识别时返回 ``None``。"""
    text = text.strip()
    if not text or any(ch not in _CN_DIGITS and ch != '十' for ch in text):
        return None

    if '十' not in text:
        value = 0
        for ch in text:
            value = value * 10 + _CN_DIGITS[ch]
        return value

    head, _, tail = text.partition('十')
    tens = _CN_DIGITS[head] if head else 1
    ones = _CN_DIGITS[tail] if tail else 0
    return tens * 10 + ones


def parse_page_number(text: str):
    """识别完整的页码标记，无法识别时返回 ``None``。

    只有整段文字符合页码格式才算命中，避免将“详见第三页”等正文误判为分页标记。
    """
    raw = str(text or '').strip()
    if not raw or len(raw) > 24:
        return None

    for pattern in (r'^第\s*(\d+)\s*页$', r'^(\d+)\s*[.、)）]?$',
                    r'^[Pp]age\s*(\d+)$', r'^[Ss]lide\s*(\d+)$'):
        match = re.match(pattern, raw)
        if match:
            return int(match.group(1))

    match = re.match(r'^第\s*([零〇一二两三四五六七八九十]+)\s*页$', raw)
    if match:
        return parse_chinese_number(match.group(1))

    return None


# PowerPoint

def parse_pptx(path: Path) -> ScriptData:
    """按幻灯片顺序读取 PowerPoint 备注。"""
    from pptx import Presentation

    presentation = Presentation(str(path))
    notes = {}
    for index, slide in enumerate(presentation.slides, start=1):
        text = ''
        for note in slide.notes_slide.notes_text_frame.paragraphs:
            text += note.text + '\n'
        notes[index] = text

    filled = sum(1 for value in notes.values() if value.strip())
    return ScriptData(
        notes=notes,
        source_name=Path(path).stem,
        report=[f'共 {len(notes)} 页，其中 {filled} 页含备注'],
    )


# Word

def _parse_docx_table(document) -> ScriptData:
    """从 Word 表格的前两列读取页码和正文。

    单元格边界定义结构，用户可自由修改正文而不会影响分页识别。
    """
    notes = {}
    report = []
    skipped = 0

    for table in document.tables:
        for row in table.rows:
            cells = row.cells
            if len(cells) < 2:
                continue

            page = parse_page_number(cells[0].text)
            if page is None:
                # 表头和页码损坏的行都不构成讲稿页；非空项计入诊断报告。
                if cells[0].text.strip():
                    skipped += 1
                continue

            text = '\n'.join(p.text for p in cells[1].paragraphs).strip()
            if page in notes:
                notes[page] = f'{notes[page]}\n{text}'.strip()
            else:
                notes[page] = text

    if notes:
        report.append(f'按表格解析，共 {len(notes)} 页')
        if skipped > 1:
            report.append(f'{skipped - 1} 行页码无法识别，已跳过')
    return ScriptData(notes=notes, report=report)


def _parse_docx_text(document) -> ScriptData:
    """按独立页码段落解析没有讲稿表格的 Word 文档。"""
    notes = {}
    report = []
    current = None
    orphan = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        page = parse_page_number(text)
        # 标题样式表达了结构边界；没有显式页码时按出现顺序分配下一页。
        is_heading = str(paragraph.style.name or '').lower().startswith('heading')
        if page is None and is_heading:
            page = (max(notes) if notes else 0) + 1

        if page is not None:
            current = page
            notes.setdefault(current, '')
            continue

        if current is None:
            orphan.append(text)
        else:
            notes[current] = f'{notes[current]}\n{text}'.strip()

    if orphan:
        first = min(notes) if notes else 1
        merged = '\n'.join(orphan)
        notes[first] = f'{merged}\n{notes.get(first, "")}'.strip()
        report.append(f'首个页码标记之前的 {len(orphan)} 段内容已归入第 {first} 页')

    if notes:
        report.insert(0, f'按段落文本解析，共 {len(notes)} 页')
    return ScriptData(notes=notes, report=report)


def parse_docx(path: Path) -> ScriptData:
    """解析 Word 讲稿，优先使用稳定的表格结构，再回退到段落结构。"""
    from docx import Document

    document = Document(str(path))

    data = _parse_docx_table(document)
    if not data.notes:
        data = _parse_docx_text(document)
        if data.notes:
            data.report.append('未找到讲稿表格，已按段落文本解析。建议重新导出以获得更稳定的往返编辑')

    if not data.notes:
        raise RuntimeError('未能识别出页码。请确认文档包含讲稿表格，'
                           '或每页以单独成行的页码标记开头')

    data.source_name = Path(path).stem
    return data
# JSON

def parse_json(path: Path) -> ScriptData:
    """解析讲稿 JSON，并用文件声明的分隔符还原逐页正文。"""
    try:
        payload = json.loads(Path(path).read_text(encoding='utf-8'))
    except Exception as e:
        raise RuntimeError(f'JSON 解析失败：{e}') from e

    if not isinstance(payload, dict):
        raise RuntimeError('JSON 顶层结构必须是对象')

    fmt = str(payload.get('format', ''))
    if fmt and fmt != JSON_FORMAT:
        raise RuntimeError(f'该文件不是本软件导出的讲稿（format={fmt}）')

    mark = str(payload.get('mark', '') or '')
    pages = payload.get('pages')
    notes = {}
    report = []

    if isinstance(pages, list):
        for entry in pages:
            if not isinstance(entry, dict):
                continue
            page = entry.get('page')
            try:
                page = int(page)
            except (TypeError, ValueError):
                report.append(f'跳过页码无法解析的条目：{page!r}')
                continue

            segments = entry.get('segments')
            if isinstance(segments, list):
                notes[page] = (mark or '●').join(str(s) for s in segments)
            else:
                notes[page] = str(entry.get('text', ''))

    if not notes:
        raise RuntimeError('JSON 中没有可用的讲稿内容')

    report.insert(0, f'按 JSON 解析，共 {len(notes)} 页')
    return ScriptData(
        notes=notes,
        source_name=str(payload.get('source_name', '') or Path(path).stem),
        mark=mark,
        report=report,
    )


# 统一入口

_PARSERS = {'.pptx': parse_pptx, '.docx': parse_docx, '.json': parse_json}


def load_script(path) -> ScriptData:
    """根据文件扩展名选择解析器，并返回统一的讲稿数据。"""
    file_path = Path(path)
    parser = _PARSERS.get(file_path.suffix.lower())
    if parser is None:
        raise RuntimeError(f'不支持的讲稿格式：{file_path.suffix}')
    return parser(file_path)
