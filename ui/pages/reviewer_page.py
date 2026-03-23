"""主页：负责导入、分段、转音频和播放"""

import json
import re
from datetime import datetime
from pathlib import Path

import pyautogui
from docx import Document
from pptx import Presentation
from PySide6.QtCore import QTimer, Qt, QUrl
from PySide6.QtGui import QIcon
from PySide6.QtMultimedia import QAudioOutput, QMediaDevices, QMediaPlayer
from PySide6.QtWidgets import QFileDialog, QWidget
from qfluentwidgets import (
    Action,
    FluentIcon,
    InfoBar,
    InfoBarPosition,
    InfoLevel,
    MessageBox,
    RoundMenu,
    setThemeColor,
)

from Ui_mainwindow import Ui_mainwindow
from app.app_context import AppContext
from tasks.audio_generation_task import AudioGenerationTask
from ui.dialogs.edit_mark_dialog import EditMarkMessageBox
from ui.dialogs.session_history_dialog import SessionHistoryDialog


class PPTReviewer(QWidget, Ui_mainwindow):
    """主页主逻辑"""

    def __init__(self, context: AppContext, parent=None):
        super().__init__(parent=parent)
        self.ctx = context
        self.setupUi(self)

        self.player = QMediaPlayer()
        self.audio_output = QAudioOutput()
        self.player.setAudioOutput(self.audio_output)
        self.player.mediaStatusChanged.connect(self.media_status_changed)

        # 监听音频设备物理插拔
        self._media_devices = QMediaDevices(self)
        self._media_devices.audioOutputsChanged.connect(self.handle_audio_device_change)

        # 轮询监听 Windows 系统“软”切换默认输出设备
        self.current_device_id = QMediaDevices.defaultAudioOutput().id()
        self.device_check_timer = QTimer(self)
        self.device_check_timer.timeout.connect(self.check_default_audio_device)
        self.device_check_timer.start(1000)

        self.media_list = []  # 存储音频文件的列表
        self.current_index = 0  # 当前播放的音频文件索引
        self.wait_current_index = 0  # 倒计时索引
        self.note_file_path = None
        self.note_file_name = ''
        self.notes = {}  # 每页讲稿
        self.notes_list = []  # 每块讲稿
        self.notes_duration_list = []
        self.note_cache_keys = []
        self.note_cache_exts = []
        self.cache_hit_count = 0  # 新增：缓存命中数
        self.is_play_notes = False
        self.is_import = False
        self.mark = '●'

        self.partingIconWidget.setIcon(QIcon(':/image/image/parting.svg'))
        self.fileIconWidget.setIcon(QIcon(':/image/image/Folder.svg'))
        self.currentIconWidget.setIcon(QIcon(':/image/image/countdown.svg'))
        self.currentTimeIconWidget.setIcon(QIcon(':/image/image/Clock.svg'))
        self.pageJumpToolButton.setIcon(FluentIcon.ACCEPT_MEDIUM)
        self.playButton.setIcon(QIcon(':/image/image/play.svg'))
        self.stopButton.setIcon(QIcon(':/image/image/stop.svg'))
        self.resetButton.setIcon(QIcon(':/image/image/backward.svg'))
        self.IndeterminateProgressBar.setVisible(False)

        self.getFileButton.setIcon(QIcon(':/image/image/ppt.svg'))
        self.file_button_menu = RoundMenu(parent=self)
        self.file_button_menu.addAction(
            Action(QIcon(':/image/image/word.svg'), '导入 Word', triggered=self.init_word_play)
        )
        self.file_button_menu.addAction(
            Action(QIcon(':/image/image/update.svg'), '历史记录列表', triggered=self.show_session_history_dialog)
        )
        self.getFileButton.setFlyout(self.file_button_menu)
        self.bgScrollArea.enableTransparentBackground()

        self.playButton.clicked.connect(self.init_play)
        self.stopButton.clicked.connect(self.stop_audio)
        self.resetButton.clicked.connect(self.reset_audio)
        self.getFileButton.clicked.connect(self.init_ppt_play)
        self.editMarkPushButton.clicked.connect(self.show_edit_mark_dialog)
        self.pageJumpToolButton.clicked.connect(self.jump_page)
        self.infoPushButton.clicked.connect(self.show_info_dialog)

        self.next_timer = QTimer(self)
        self.next_timer.timeout.connect(self.timeout_play_next_audio)

        self.wav_temp_path = Path('./temp').resolve()
        self.countdown_wav_temp_path = Path('./data/cache/countdown').resolve()
        self.audio_cache_path = Path('./data/cache/audio_chunks').resolve()
        self.session_root_path = Path('./data/sessions').resolve()
        self.wav_temp_path.mkdir(parents=True, exist_ok=True)
        self.countdown_wav_temp_path.mkdir(parents=True, exist_ok=True)
        self.audio_cache_path.mkdir(parents=True, exist_ok=True)
        self.session_root_path.mkdir(parents=True, exist_ok=True)

        self.save_thread = AudioGenerationTask(self, self.ctx.tts_engine, self)
        self.save_thread.signal_import_index.connect(self.thread_print_index)
        self.save_thread.signal_cache_hit_count.connect(self.set_cache_hit_count)
        self.save_thread.signal_finish.connect(self.thread_save_finish)

        self.check_import()

    def check_import(self):
        """检查导入状态"""
        if self.is_import:
            self.playCardWidget.setEnabled(True)
            self.playCardWidget_2.setEnabled(True)
            self.playCardWidget_3.setEnabled(True)
            self.statusLabel.setText('已导入')
            self.IconInfoBadge.setLevel(InfoLevel.SUCCESS)
            self.IconInfoBadge.setIcon(FluentIcon.ACCEPT_MEDIUM)
        else:
            self.playCardWidget.setEnabled(False)
            self.playCardWidget_2.setEnabled(False)
            self.playCardWidget_3.setEnabled(False)
            self.statusLabel.setText('未导入')
            self.IconInfoBadge.setLevel(InfoLevel.INFOAMTION)
            self.IconInfoBadge.setIcon(FluentIcon.ACCEPT_MEDIUM)

    def handle_audio_device_change(self):
        """处理设备物理插拔后的自适应切换"""
        self.check_default_audio_device()

    def check_default_audio_device(self):
        """发现默认输出设备变化时自动切换"""
        default_device = QMediaDevices.defaultAudioOutput()
        if default_device.id() != self.current_device_id:
            self.current_device_id = default_device.id()
            self.audio_output.setDevice(default_device)
            if self.player.playbackState() == QMediaPlayer.PlaybackState.PlayingState:
                pos = self.player.position()
                self.player.pause()
                self.player.play()
                self.player.setPosition(pos)
            print(f'[播放器] 音频输出设备已自适应切换至: {default_device.description()}')

    def play_audio(self):
        """播放音频文件"""
        if self.is_play_notes:
            if self.current_index < len(self.media_list):
                self.player.setSource(QUrl.fromLocalFile(str(self.media_list[self.current_index])))
                self.player.play()
                self.playButton.setEnabled(False)
                self.currentStatusLabel.setText('播放')
                self.set_current_label_text()
                print(self.notes_list[self.current_index]['text'])
            else:
                print('播放完毕')
                self.reset_audio()
        else:
            if self.wait_current_index < len(self.media_list):
                self.player.setSource(QUrl.fromLocalFile(str(self.media_list[self.wait_current_index])))
                self.player.play()
                self.playButton.setEnabled(False)
                temp_index = len(self.media_list) - self.wait_current_index
                self.currentStatusLabel.setText('倒计时')
                self.currentPageLabel.setText(f'{temp_index}')
                self.currentIndexLabel.setText(f'{temp_index}')
            else:
                print('播放完毕')
                self.wait_current_index = 0
                self.play_notes()

    def media_status_changed(self, status):
        """媒体状态转移"""
        if status == QMediaPlayer.MediaStatus.EndOfMedia:
            self.next_timer.start(100)

    def stop_audio(self):
        """停止播放"""
        self.player.stop()
        self.playButton.setEnabled(True)
        self.currentStatusLabel.setText('停止')

        try:
            if self.window() and hasattr(self.window(), 'setting_interface'):
                if hasattr(self.window().setting_interface, 'preview_player'):
                    self.window().setting_interface.preview_player.stop()
        except Exception:
            pass

    def reset_audio(self):
        """重置播放"""
        self.stop_audio()
        self.current_index = 0
        if self.notes_list:
            self.set_current_label_text()

    def set_current_label_text(self):
        """更新当前索引标签"""
        self.currentPageLabel.setText(f'{self.notes_list[self.current_index]["page"]} / {len(self.notes)}')
        self.currentIndexLabel.setText(f'{self.current_index + 1} / {len(self.notes_list)}')

    def get_index_from_page(self, page):
        """根据页码查找首条讲稿索引"""
        for i, item in enumerate(self.notes_list):
            if item['page'] == page:
                return i
        return -1

    def timeout_play_next_audio(self):
        """当前音频结束后自动播放下一个"""
        self.next_timer.stop()
        if self.is_play_notes:
            if self.scrollEnableSwitch.isChecked():
                pyautogui.press('pagedown')
            self.current_index += 1
        else:
            self.wait_current_index += 1
        self.play_audio()

    def init_word_play(self):
        """初始化 Word 导入"""
        self.getFileButton.setEnabled(False)
        self.get_word_path()
        if not self.note_file_path:
            print('文件未导入！')
            self.create_warning_info_bar('导入已取消', '请重新选择文件进行导入。')
            self.getFileButton.setEnabled(True)
            return False
        self.notesPathLabel.setText(str(self.note_file_path))

        setThemeColor('#2B579A')

        try:
            self.get_word_notes_dict()
        except Exception as e:
            print(e)
            self.create_error_info_bar('Word 文件解析错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False
        self.init_general_play()
        return True

    def init_ppt_play(self):
        """初始化 PPT 导入"""
        self.getFileButton.setEnabled(False)
        self.get_ppt_path()
        if not self.note_file_path:
            print('文件未导入！')
            self.create_warning_info_bar('导入已取消', '请重新选择文件进行导入')
            self.getFileButton.setEnabled(True)
            return False
        self.notesPathLabel.setText(str(self.note_file_path))

        setThemeColor('#B7472A')

        try:
            self.get_ppt_notes_dict()
        except Exception as e:
            print(e)
            self.create_error_info_bar('PowerPoint 文件解析错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False
        self.init_general_play()
        return True

    def init_general_play(self):
        """通用初始化流程"""
        try:
            self.mark_split()
        except Exception as e:
            print(e)
            self.create_error_info_bar('讲稿解析错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False

        self.clean_and_reset()

        try:
            self.clean_temp_folder(self.wav_temp_path)
        except Exception as e:
            print(e)
            self.create_error_info_bar('缓存清理错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False

        try:
            self.save_thread.start()
            self.IndeterminateProgressBar.setVisible(True)
        except Exception as e:
            print(e)
            self.create_error_info_bar('语音转换错误', f'详情：{e}')
            self.getFileButton.setEnabled(True)
            return False

    def clean_and_reset(self):
        """清理已导入内容并重置状态"""
        self.stop_audio()
        self.player.setSource(QUrl())

        try:
            if self.window() and hasattr(self.window(), 'setting_interface'):
                if hasattr(self.window().setting_interface, 'preview_player'):
                    self.window().setting_interface.preview_player.stop()
                    self.window().setting_interface.preview_player.setSource(QUrl())
        except Exception:
            pass

        self.media_list = []
        self.current_index = 0
        self.notes_duration_list = []
        self.note_cache_keys = []
        self.note_cache_exts = []
        self.cache_hit_count = 0
        self.is_import = False
        self.check_import()

    def refresh_notes_duration_list(self):
        """按当前已生成的音频文件重新统计时长"""
        if not self.media_list or len(self.media_list) != len(self.notes_list):
            self.load_audio_files()
        duration_list = []
        for path in self.media_list:
            duration = AudioGenerationTask.get_audio_duration(path)
            duration_list.append(duration)
        self.notes_duration_list = duration_list

    def get_word_path(self):
        """获取 Word 路径"""
        selected_files = QFileDialog.getOpenFileName(self, '选择Word文件', '', 'Word Files (*.docx)')
        self.note_file_path = Path(selected_files[0]) if selected_files[0] else None
        self.set_filename()

    def get_ppt_path(self):
        """获取 PPT 路径"""
        selected_files = QFileDialog.getOpenFileName(self, '选择PowerPoint文件', '', 'PowerPoint Files (*.pptx)')
        self.note_file_path = Path(selected_files[0]) if selected_files[0] else None
        self.set_filename()

    def set_filename(self):
        """根据文件路径生成文件名"""
        if not self.note_file_path:
            self.note_file_name = ''
            return
        self.note_file_name = self.note_file_path.stem

    def get_word_notes_dict(self):
        """从 Word 获取讲稿字典"""
        doc = Document(str(self.note_file_path))
        page_data = {}
        current_page = None
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text.startswith('第') and text.endswith('页'):
                page_text = text[1:-1]
                try:
                    current_page = int(page_text)
                except ValueError:
                    current_page = None
            elif current_page is not None:
                if current_page in page_data:
                    page_data[current_page] += '\n' + text
                else:
                    page_data[current_page] = text
        self.notes = page_data

    def get_ppt_notes_dict(self):
        """从 PPT 备注获取讲稿字典"""
        presentation = Presentation(str(self.note_file_path))
        notes_dict = {}
        for i, slide in enumerate(presentation.slides, start=1):
            text = ''
            for note in slide.notes_slide.notes_text_frame.paragraphs:
                text += note.text + '\n'
            notes_dict[i] = text
        self.notes = notes_dict

    def mark_split(self):
        """按分隔符分割讲稿"""
        notes_list = []
        for page in range(1, len(self.notes) + 1):
            note_text = re.sub(r'[\x00-\x1f\x7f]+', '', self.notes[page])
            if note_text and note_text[-1] == self.mark:
                note_text = note_text[:-1]
            note_list = note_text.split(self.mark)
            for one_note in note_list:
                notes_list.append({'page': page, 'text': one_note.strip()})
        self.notes_list = notes_list
        print('讲稿分割完毕')

    @staticmethod
    def clean_temp_folder(path: Path):
        """清理临时音频（wav/mp3）"""
        for pattern in ('*.wav', '*.mp3'):
            for file_path in path.glob(pattern):
                try:
                    file_path.unlink()
                    print(f'已清理 {file_path.name}')
                except Exception as e:
                    print(f'清理文件失败: {file_path.name}, 原因: {e}')
        print('转换完成')

    def thread_print_index(self, import_index):
        """生成进度回调"""
        text = f'已生成：{import_index}/{len(self.notes_list)}'
        print(text)
        self.statusLabel.setText(text)

    def set_cache_hit_count(self, count: int):
        """设置缓存命中数"""
        self.cache_hit_count = count

    def thread_save_finish(self):
        """音频生成完成回调"""
        print('转换完成')
        self.refresh_notes_duration_list()
        self.save_session_record()
        self.IndeterminateProgressBar.setVisible(False)
        
        # 构建提示文本，包含缓存命中信息
        total_notes = len(self.notes_list)
        if self.cache_hit_count > 0:
            cache_info = f'（命中缓存 {self.cache_hit_count}/{total_notes} 条）'
            self.create_success_info_bar('转换完成', f'音频播放功能已准备就绪 {cache_info}')
        else:
            self.create_success_info_bar('转换完成', '音频播放功能已准备就绪')
        self.getFileButton.setEnabled(True)
        self.is_import = True
        self.check_import()
        if self.notes_list:
            self.set_current_label_text()
        self.pageJumpSpinBox.setMaximum(len(self.notes))

    def save_session_record(self):
        """保存当前导入会话，供后续快速恢复"""
        if not self.notes_list:
            return

        generation_profile = self.ctx.tts_engine.get_generation_profile()
        if len(self.note_cache_keys) != len(self.notes_list):
            self.note_cache_keys = [
                self.ctx.tts_engine.build_audio_cache_key(item['text'], generation_profile)
                for item in self.notes_list
            ]

        if len(self.note_cache_exts) != len(self.notes_list):
            output_ext = self.ctx.tts_engine.get_output_extension()
            self.note_cache_exts = [output_ext] * len(self.notes_list)

        durations = self.notes_duration_list[:]
        if len(durations) != len(self.notes_list):
            durations = [0.0] * len(self.notes_list)

        items = []
        for index, note in enumerate(self.notes_list):
            items.append({
                'index': index,
                'page': note['page'],
                'text': note['text'],
                'duration': float(durations[index]),
                'cache_key': self.note_cache_keys[index],
                'cache_ext': self.note_cache_exts[index],
            })

        now = datetime.now()
        session_id = now.strftime('%Y%m%d_%H%M%S')
        record = {
            'version': 1,
            'session_id': session_id,
            'created_at': now.isoformat(timespec='seconds'),
            'source_file': str(self.note_file_path) if self.note_file_path else '',
            'source_name': self.note_file_name,
            'mark': self.mark,
            'generation_profile': generation_profile,
            'notes': self.notes,
            'items': items,
        }

        record_path = self.session_root_path / f'{session_id}.json'
        with record_path.open('w', encoding='utf-8') as f:
            json.dump(record, f, ensure_ascii=False, indent=2)

    def load_session_record_from_file(self):
        """从历史记录文件恢复导入会话"""
        selected = QFileDialog.getOpenFileName(
            self,
            '选择历史记录文件',
            str(self.session_root_path),
            'Session Files (*.json)'
        )
        if not selected[0]:
            return

        try:
            self.load_session_record(Path(selected[0]))
        except Exception as e:
            print(e)
            self.create_error_info_bar('历史记录加载失败', f'详情：{e}')

    def show_session_history_dialog(self):
        """显示历史记录列表并加载选中项"""
        dialog = SessionHistoryDialog(self.session_root_path, self)
        if not dialog.exec():
            return

        record_path = dialog.get_selected_record_path()
        if not record_path:
            self.create_warning_info_bar('未选择记录', '请在历史记录列表中选择一条记录。')
            return

        try:
            self.load_session_record(record_path)
        except Exception as e:
            print(e)
            self.create_error_info_bar('历史记录加载失败', f'详情：{e}')

    def load_session_record(self, record_path: Path):
        """加载指定会话记录并恢复音频播放状态"""
        with record_path.open('r', encoding='utf-8') as f:
            record = json.load(f)

        items = record.get('items', [])
        if not items:
            raise RuntimeError('历史记录为空，无法加载')

        self.clean_and_reset()

        missing_list = []
        media_list = []
        notes_list = []
        duration_list = []
        cache_keys = []
        cache_exts = []

        for idx, item in enumerate(items):
            page = int(item.get('page', 0))
            text = str(item.get('text', ''))
            cache_key = str(item.get('cache_key', '')).strip()
            if not cache_key:
                profile = record.get('generation_profile', {})
                cache_key = self.ctx.tts_engine.build_audio_cache_key(text, profile)

            cache_ext = str(item.get('cache_ext', '')).strip().lower().lstrip('.')
            ext_candidates = [cache_ext] if cache_ext else []
            for ext in ('wav', 'mp3'):
                if ext not in ext_candidates:
                    ext_candidates.append(ext)

            cache_path = None
            for ext in ext_candidates:
                candidate = self.audio_cache_path / f'{cache_key}.{ext}'
                if candidate.exists() and candidate.stat().st_size > 0:
                    cache_path = candidate
                    cache_ext = ext
                    break

            if cache_path is None:
                missing_list.append(f'第{page}页-第{idx + 1}条')
                continue

            media_list.append(cache_path)
            notes_list.append({'page': page, 'text': text})
            duration_list.append(float(item.get('duration', 0.0)))
            cache_keys.append(cache_key)
            cache_exts.append(cache_ext)

        if missing_list:
            missing_text = '、'.join(missing_list[:10])
            raise RuntimeError(f'存在缺失音频缓存：{missing_text}')

        self.note_file_name = str(record.get('source_name', ''))
        source_file = str(record.get('source_file', '')).strip()
        self.note_file_path = Path(source_file) if source_file else None
        self.mark = str(record.get('mark', self.mark))

        notes = record.get('notes')
        if isinstance(notes, dict) and notes:
            self.notes = {int(k): v for k, v in notes.items()}
        else:
            rebuilt_notes = {}
            for item in notes_list:
                rebuilt_notes.setdefault(item['page'], []).append(item['text'])
            self.notes = {k: self.mark.join(v) for k, v in rebuilt_notes.items()}

        self.notes_list = notes_list
        self.notes_duration_list = duration_list
        self.note_cache_keys = cache_keys
        self.note_cache_exts = cache_exts

        self.media_list = media_list
        self.current_index = 0
        self.is_import = True
        self.check_import()
        self.pageJumpSpinBox.setMaximum(len(self.notes))
        if self.notes_list:
            self.set_current_label_text()

        self.notesPathLabel.setText(f'历史记录：{record_path.name}')
        self.create_success_info_bar('加载成功', '历史记录已恢复，可直接播放')


    def init_play(self):
        """准备播放"""
        if self.currentSwitch.isChecked():
            self.play_wait()
        else:
            self.play_notes()

    def play_wait(self):
        """播放倒计时"""
        self.is_play_notes = False
        self.wait_current_index = 0
        self.load_wait_audio_files()
        print('已导入倒计时')
        self.play_audio()

    def play_notes(self):
        """播放讲稿"""
        self.is_play_notes = True
        if not self.media_list or 'countdown' in str(self.media_list[0]):
            self.load_audio_files()
        self.play_audio()

    def load_audio_files(self):
        """查找所有正文音频（wav/mp3），添加到 media_list 中"""
        audio_files = list(self.wav_temp_path.glob('*.wav')) + list(self.wav_temp_path.glob('*.mp3'))
        audio_files = sorted(
            audio_files,
            key=lambda path: [int(part) if part.isdigit() else part for part in path.stem.split('_')]
        )
        self.media_list = audio_files
        print('音频列表载入完成')

    def load_wait_audio_files(self):
        """获取指定数量的倒计时 wav，添加到 media_list 中"""
        audio_files = [
            path for path in self.countdown_wav_temp_path.glob('*.wav')
            if path.stem.isdigit() and int(path.stem) <= self.currentSpinBox.value()
        ]
        audio_files = sorted(audio_files, key=lambda path: int(path.stem), reverse=True)
        self.media_list = audio_files
        print('倒计时列表载入完成')

    def jump_page(self):
        """跳转到指定页"""
        self.stop_audio()
        index = self.get_index_from_page(self.pageJumpSpinBox.value())
        if index > -1:
            self.current_index = index
        if self.notes_list:
            self.set_current_label_text()

    def create_success_info_bar(self, title, text):
        """成功消息框"""
        InfoBar.success(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=2000,
            parent=self
        )

    def create_warning_info_bar(self, title, text):
        """警告消息框"""
        InfoBar.warning(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=3000,
            parent=self
        )

    def create_error_info_bar(self, title, text):
        """错误消息框"""
        InfoBar.error(
            title=title,
            content=text,
            orient=Qt.Horizontal,
            isClosable=True,
            position=InfoBarPosition.TOP,
            duration=-1,
            parent=self
        )

    def show_edit_mark_dialog(self):
        """编辑分隔符弹窗"""
        box = EditMarkMessageBox(self.mark, self)
        if box.exec():
            text = box.urlLineEdit.text()
            print(f'分隔符：{text}')
            self.mark = text

    @staticmethod
    def s_to_str(s):
        """将秒转换为时间标签"""
        if s < 60:
            return f'{round(s, 2)} 秒'
        m, s = divmod(s, 60)
        if m < 60:
            return f'{round(m)} 分钟 {round(s)} 秒'
        h, m = divmod(m, 60)
        return f'{round(h)} 小时 {round(m)} 分钟 {round(s)} 秒'

    def count_words(self):
        """统计总字数"""
        text = ''
        for page in self.notes:
            text += self.notes[page]
        text = re.sub(r'\s+', '', text)
        return len(text)

    def show_info_dialog(self):
        """显示统计弹窗"""
        title = '统计信息'
        if len(self.notes_duration_list) != len(self.notes_list):
            self.refresh_notes_duration_list()
        if not self.notes_duration_list or not self.notes_list:
            self.create_warning_info_bar('暂无统计信息', '请先导入并生成音频后再查看统计。')
            return

        words_count = self.count_words()
        max_duration = max(self.notes_duration_list)
        max_duration_index = self.notes_duration_list.index(max_duration)
        min_duration = min(self.notes_duration_list)
        min_duration_index = self.notes_duration_list.index(min_duration)

        content_list = [
            ['页码总计', f'{len(self.notes)} 页'],
            ['音频总计', f'{len(self.notes_list)} 条'],
            ['演讲稿字数总计', f'{words_count} 字'],
            ['音频总时长', f'{self.s_to_str(sum(self.notes_duration_list))}\n'],
            ['最长音频时长', f'{self.s_to_str(max_duration)}'],
            ['最长音频索引', f'{max_duration_index}'],
            ['最长音频所属页码', f'第 {self.notes_list[max_duration_index]["page"]} 页'],
            ['最短音频时长', f'{self.s_to_str(min_duration)}'],
            ['最短音频索引', f'{min_duration_index}'],
            ['最短音频所属页码', f'第 {self.notes_list[min_duration_index]["page"]} 页'],
        ]

        content = ''
        for item in content_list:
            content += f'{item[0]}：{item[1]}\n'
        dialog = MessageBox(title, content, self)
        dialog.yesButton.setText('确定')
        dialog.cancelButton.setVisible(False)
        dialog.exec()
