"""提供 AI 讲稿处理以及文档、字幕、音频和工程包导出。"""

import json
from pathlib import Path

from PySide6.QtCore import QSize, Qt
from PySide6.QtWidgets import QFileDialog, QHBoxLayout, QVBoxLayout, QWidget
from qfluentwidgets import (
    CaptionLabel,
    CardWidget,
    InfoBar,
    InfoBarPosition,
    MessageBox,
    PrimaryPushButton,
    PushButton,
    SubtitleLabel,
)

from app import audio_cache, exporters, project_package, theme
from app.app_context import AppContext
from toolsInterface import Ui_toolsInterface
from ui.dialogs.rewrite_dialog import RewriteDialog


class ToolsInterface(QWidget, Ui_toolsInterface):
    """根据主页当前会话提供讲稿处理与多格式导出操作。"""

    def __init__(self, context: AppContext, reviewer_page, parent=None):
        super().__init__(parent=parent)
        self.ctx = context
        self.reviewer_page = reviewer_page
        self.setupUi(self)
        self.bgScrollArea.enableTransparentBackground()

        self.processing_cards = []
        self.export_cards = []
        self.build_layout()
        self.refresh_state()

    # 页面装配

    def add_section(self, title: str, caption: str = ''):
        """向工具页追加分区标题和可选说明。"""
        label = SubtitleLabel(title, self.importWidget)
        label.setMinimumSize(QSize(0, 30))
        label.setMaximumSize(QSize(16777215, 30))
        self.verticalLayout_2.addWidget(label)

        caption_label = None
        if caption:
            caption_label = CaptionLabel(caption, self.importWidget)
            caption_label.setWordWrap(True)
            self.verticalLayout_2.addWidget(caption_label)
        return label, caption_label

    def add_card(self, title: str, caption: str, controls,
                 primary_first: bool = False, button_width: int = 120):
        """追加操作卡片；``controls`` 提供按钮文本与槽函数。"""
        card = CardWidget(self.importWidget)
        h_layout = QHBoxLayout(card)
        h_layout.setContentsMargins(20, 16, 20, 16)

        info_layout = QVBoxLayout()
        info_layout.setSpacing(0)

        title_label = theme.make_card_title(title, card)
        info_layout.addWidget(title_label)

        caption_label = CaptionLabel(card)
        caption_label.setText(caption)
        caption_label.setWordWrap(True)
        info_layout.addWidget(caption_label)

        h_layout.addLayout(info_layout, stretch=1)
        h_layout.addSpacing(16)

        for index, (text, slot) in enumerate(controls):
            factory = PrimaryPushButton if (primary_first and index == 0) else PushButton
            button = factory(text, card)
            button.clicked.connect(slot)
            button.setMinimumSize(QSize(button_width, 33))
            button.setMaximumSize(QSize(button_width, 16777215))
            h_layout.addWidget(button)

        self.verticalLayout_2.addWidget(card)
        return card

    def build_layout(self):
        """将 Designer 骨架重排为讲稿处理和导出两个分区。"""
        # 复用 .ui 中的静态标题，避免运行期重建造成样式差异。
        self.SubtitleLabel.setText('讲稿处理')
        self.CaptionLabel_4.setText('对已导入的讲稿进行加工')

        self.processing_cards.append(self.add_card(
            'AI 优化',
            '通过 OpenAI 兼容接口逐页优化讲稿',
            [('前往优化', self.open_rewrite_dialog)],
            primary_first=True,
        ))

        self.add_section('导出')

        self.export_cards.append(self.add_card(
            '讲稿文档',
            '写回 PowerPoint 备注，或导出为其他格式的讲稿文档',
            [
                ('PPT 备注', self.write_to_ppt),
                ('Word', self.write_to_word),
                ('Markdown', self.write_to_markdown),
                ('JSON', self.write_to_json),
            ],
        ))

        self.export_cards.append(self.add_card(
            '工程包',
            '将讲稿、音频与生成配置打包为单个文件，可在其它设备导入播放',
            [('导出工程包', self.export_project_package)],
        ))

        self.export_cards.append(self.add_card(
            '字幕与音频',
            '导出字幕文件与已合成的语音音频，可选择逐条或合并为单个文件',
            [
                ('SRT 字幕', self.write_to_srt),
                ('逐条音频', self.export_audio_files),
                ('合并音频', self.export_merged_audio),
            ],
        ))


    # 状态同步

    def showEvent(self, event):
        """页面显示时从主页同步最新会话状态。"""
        super().showEvent(event)
        self.refresh_state()

    def refresh_state(self):
        """按讲稿是否就绪更新各操作按钮。"""
        imported = bool(self.reviewer_page.notes)
        for card in self.processing_cards + self.export_cards:
            # 仅禁用操作按钮，保留卡片说明的正常对比度。
            for button in card.findChildren(PushButton):
                button.setEnabled(imported)

    # 讲稿处理

    def open_rewrite_dialog(self):
        """打开逐页改写对照框，并将用户确认的结果写回主页。"""
        if not self.require_notes() or not self.require_idle():
            return

        page = self.reviewer_page
        dialog = RewriteDialog(page.notes, page.mark, self.ctx.app_settings, self)

        accepted = dialog.exec()
        # 弹窗中可能只修改了接口或风格，即使未采用结果也应保存这些设置。
        self.ctx.config.save_later()
        if not accepted:
            return

        applied = page.apply_rewritten_notes(dialog.get_applied())
        if applied <= 0:
            self.create_warning_info_bar('未写回任何内容', '改写结果为空或页码不匹配。')
            return

        self.goto_reviewer(f'已更新 {applied} 页讲稿，正在重新合成受影响的音频')

    def goto_reviewer(self, message: str):
        """返回主页，并展示后续音频生成状态。"""
        window = self.window()
        if hasattr(window, 'switchTo'):
            window.switchTo(self.reviewer_page)
        self.reviewer_page.create_success_info_bar('已开始重新合成', message)

    # 共用前置检查

    def require_notes(self) -> bool:
        """检查讲稿是否已导入，未就绪时在当前页提示。"""
        if not self.reviewer_page.notes:
            self.create_warning_info_bar('演讲稿未导入', '请先在主页导入演讲稿。')
            return False
        return True

    def require_idle(self) -> bool:
        """检查音频生成任务是否空闲。"""
        if self.reviewer_page.is_busy():
            self.create_warning_info_bar('正在生成音频', '请等待当前转换完成后再试。')
            return False
        return True

    def require_durations(self) -> bool:
        """检查所有讲稿段是否已有对应的时长数据。"""
        page = self.reviewer_page
        if not page.notes_list or len(page.notes_duration_list) != len(page.notes_list):
            self.create_warning_info_bar('音频尚未就绪', '请先完成语音转换后再导出。')
            return False
        return True

    # 导出目录

    def choose_export_dir(self):
        """让用户选择目标目录，并以上次选择作为起始位置。"""
        start_dir = str(self.ctx.app_settings.get('export_dir') or '')
        selected = QFileDialog.getExistingDirectory(
            self, '选择导出目录', start_dir, QFileDialog.Option.ShowDirsOnly)
        if not selected:
            return None

        # 只保存目录偏好，不代表后续导出可跳过确认。
        self.ctx.app_settings.set('export_dir', selected)
        self.ctx.config.save_later()
        return Path(selected)

    @staticmethod
    def unique_path(path: Path) -> Path:
        """为已存在的目标文件追加递增序号，避免静默覆盖。"""
        if not path.exists():
            return path
        for index in range(2, 1000):
            candidate = path.with_name(f'{path.stem}_{index}{path.suffix}')
            if not candidate.exists():
                return candidate
        return path

    def get_ppt_path(self):
        """选择待回写的 PowerPoint，优先定位到当前源文件。"""
        current = self.reviewer_page.note_file_path
        start_path = str(current) if current and current.suffix.lower() == '.pptx' else ''
        selected_files = QFileDialog.getOpenFileName(
            self, '选择要写入备注的 PowerPoint 文件', start_path, 'PowerPoint Files (*.pptx)')
        return Path(selected_files[0]) if selected_files[0] else None

    # 讲稿文档导出

    def write_to_ppt(self):
        """将当前逐页讲稿回写到指定 PowerPoint 的备注。"""
        if not self.require_notes():
            return
        notes_dict = self.reviewer_page.notes

        ppt_path = self.get_ppt_path()
        if not ppt_path:
            self.create_warning_info_bar('已取消', '未选择要写入备注的 PowerPoint 文件。')
            return

        dir_path = self.choose_export_dir()
        if not dir_path:
            return

        try:
            from pptx import Presentation

            ppt = Presentation(str(ppt_path))
            for slide_number, slide in enumerate(ppt.slides, start=1):
                if slide_number in notes_dict:
                    slide.notes_slide.notes_text_frame.text = notes_dict[slide_number]
        except Exception as e:
            print(e)
            self.create_error_info_bar('PowerPoint 读取错误', f'详情：{e}')
            return

        output_path = self.unique_path(dir_path / f'{ppt_path.stem}_NEW.pptx')
        try:
            ppt.save(str(output_path))
        except Exception as e:
            print(e)
            self.create_error_info_bar('PowerPoint 保存错误', f'详情：{e}')
            return

        self.create_success_info_bar('生成成功', f'讲稿备注已写入：{output_path}')

    def write_to_word(self):
        """导出可修改后稳定回导的 Word 表格讲稿。"""
        if not self.require_notes():
            return

        dir_path = self.choose_export_dir()
        if not dir_path:
            return

        page = self.reviewer_page
        file_path = self.unique_path(dir_path / f'{page.note_file_name}_Notes.docx')
        try:
            exporters.write_docx(
                file_path, page.notes_list, page.notes_duration_list, mark=page.mark,
            )
        except Exception as e:
            print(e)
            self.create_error_info_bar('Word 保存错误', f'详情：{e}')
            return

        self.create_success_info_bar('转换成功', f'Word 已导出：{file_path}')

    def write_to_markdown(self):
        """导出适合阅读和打印的 Markdown 讲稿。"""
        if not self.require_notes():
            return

        dir_path = self.choose_export_dir()
        if not dir_path:
            return

        page = self.reviewer_page
        engine_def = self.ctx.tts_engine.get_current_engine_definition()
        content = exporters.build_markdown(
            page.notes_list,
            page.notes_duration_list,
            title=page.note_file_name or '演讲稿',
            mark=page.mark,
            engine_name=str(engine_def.get('name', '')),
            speaker=self.ctx.tts_engine.get_selected_voice_name(),
        )

        file_path = self.unique_path(dir_path / f'{page.note_file_name}_Notes.md')
        try:
            file_path.write_text(content, encoding='utf-8')
        except Exception as e:
            print(e)
            self.create_error_info_bar('Markdown 保存错误', f'详情：{e}')
            return

        self.create_success_info_bar('转换成功', f'Markdown 已导出：{file_path}')

    def write_to_json(self):
        """导出保留分页和分段结构的可回导 JSON。"""
        if not self.require_notes():
            return

        dir_path = self.choose_export_dir()
        if not dir_path:
            return

        page = self.reviewer_page
        engine_def = self.ctx.tts_engine.get_current_engine_definition()
        payload = exporters.build_script_json(
            page.notes_list, page.notes_duration_list,
            source_name=page.note_file_name, mark=page.mark,
            engine_name=str(engine_def.get('name', '')),
            speaker=self.ctx.tts_engine.get_selected_voice_name(),
        )

        file_path = self.unique_path(dir_path / f'{page.note_file_name}_Notes.json')
        try:
            with file_path.open('w', encoding='utf-8') as json_file:
                json.dump(payload, json_file, ensure_ascii=False, indent=2)
        except Exception as e:
            print(e)
            self.create_error_info_bar('JSON 保存错误', f'详情：{e}')
            return

        self.create_success_info_bar('转换成功', f'JSON 已导出：{file_path}')

    def export_project_package(self):
        """将讲稿、音频和生成配置打包为可跨设备恢复的工程包。"""
        page = self.reviewer_page
        if not self.require_notes():
            return
        if not page.media_list:
            self.create_warning_info_bar('音频尚未就绪', '请先完成语音转换后再导出工程包。')
            return

        items = project_package.build_items(
            page.notes_list, page.notes_duration_list,
            page.note_cache_keys, page.note_cache_exts, page.media_list)
        audio_size = project_package.estimate_size(items)

        box = MessageBox(
            '导出工程包',
            f'将打包 {len(items)} 条语句及对应音频，音频约占 '
            f'{audio_cache.format_size(audio_size)}。'
            f'\n再次导入可直接播放，无需额外配置。',
            self,
        )
        box.yesButton.setText('选择保存位置')
        box.cancelButton.setText('取消')
        if not box.exec():
            return

        dir_path = self.choose_export_dir()
        if not dir_path:
            return

        target = self.unique_path(
            dir_path / f'{page.note_file_name}{project_package.PACKAGE_SUFFIX}')
        try:
            info = project_package.export_package(
                target, items=items, notes=page.notes, mark=page.mark,
                source_name=page.note_file_name,
                generation_profile=self.ctx.tts_engine.get_generation_profile(),
                speaker=self.ctx.tts_engine.get_selected_voice_name(),
                app_version=self.ctx.version,
            )
        except Exception as e:
            print(e)
            self.create_error_info_bar('工程包导出失败', f'详情：{e}')
            return

        detail = f'已打包 {info["audio_count"]} 条音频，大小 {audio_cache.format_size(info["size"])}'
        if info['missing']:
            detail += f'；{len(info["missing"])} 条音频缺失，未包含在内'
        self.create_success_info_bar('导出成功', f'{detail}：{target}')

    # 字幕与音频导出

    def write_to_srt(self):
        """按实际音频时长导出连续 SRT 字幕。"""
        if not self.require_notes() or not self.require_durations():
            return

        dir_path = self.choose_export_dir()
        if not dir_path:
            return

        page = self.reviewer_page
        content = exporters.build_srt(page.notes_list, page.notes_duration_list)
        if not content.strip():
            self.create_warning_info_bar('没有可导出的字幕', '当前讲稿没有有效的语句时长。')
            return

        file_path = self.unique_path(dir_path / f'{page.note_file_name}_Notes.srt')
        try:
            file_path.write_text(content, encoding='utf-8')
        except Exception as e:
            print(e)
            self.create_error_info_bar('字幕保存错误', f'详情：{e}')
            return

        self.create_success_info_bar('转换成功', f'字幕已导出：{file_path}')

    def export_audio_files(self):
        """将每条讲稿音频复制为带顺序和页码的独立文件。"""
        page = self.reviewer_page
        if not page.media_list:
            self.create_warning_info_bar('音频尚未就绪', '请先完成语音转换后再导出。')
            return

        dir_path = self.choose_export_dir()
        if not dir_path:
            return

        target_dir = dir_path / f'{page.note_file_name}_Audio'
        try:
            exported = exporters.export_audio_files(page.media_list, page.notes_list, target_dir)
        except Exception as e:
            print(e)
            self.create_error_info_bar('音频导出错误', f'详情：{e}')
            return

        self.create_success_info_bar('导出成功', f'已导出 {exported} 条音频到：{target_dir}')

    def export_merged_audio(self):
        """将格式兼容的全部讲稿音频合并为单个文件。"""
        page = self.reviewer_page
        if not page.media_list:
            self.create_warning_info_bar('音频尚未就绪', '请先完成语音转换后再导出。')
            return

        dir_path = self.choose_export_dir()
        if not dir_path:
            return

        suffix = Path(page.media_list[0]).suffix.lower() or '.wav'
        file_path = self.unique_path(dir_path / f'{page.note_file_name}_Full{suffix}')
        try:
            exporters.merge_audio(page.media_list, file_path)
        except Exception as e:
            print(e)
            self.create_error_info_bar('音频合并失败', f'详情：{e}')
            return

        self.create_success_info_bar('导出成功', f'完整音频已导出：{file_path}')

    # 消息提示

    def create_success_info_bar(self, title, text):
        """在工具页顶部显示短暂的成功提示。"""
        InfoBar.success(title=title, content=text, orient=Qt.Horizontal, isClosable=True,
                        position=InfoBarPosition.TOP, duration=5000, parent=self)

    def create_warning_info_bar(self, title, text):
        """在工具页顶部显示短暂的警告提示。"""
        InfoBar.warning(title=title, content=text, orient=Qt.Horizontal, isClosable=True,
                        position=InfoBarPosition.TOP, duration=5000, parent=self)

    def create_error_info_bar(self, title, text):
        """在工具页顶部显示需要手动关闭的错误提示。"""
        InfoBar.error(title=title, content=text, orient=Qt.Horizontal, isClosable=True,
                      position=InfoBarPosition.TOP, duration=-1, parent=self)
